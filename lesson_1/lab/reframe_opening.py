"""Reframe the opening slides of L1_outline.pptx from a generic 'build AI agents'
course to THIS course: Evaluating AI Systems — Experiment Design (eval, estimation,
experiment design, uncertainty). Designs/formatting preserved; only text swapped.

Edits text in place by matching the current paragraph text (whitespace-normalized)
and rewriting the paragraph's first run, so run-level formatting is kept. Any
unmatched 'old' string is reported so nothing fails silently.
"""
import re
from pathlib import Path
from pptx import Presentation

SRC = str(Path(__file__).resolve().parent.parent / "L1_generated.pptx")
prs = Presentation(SRC)


def norm(s):
    return re.sub(r"\s+", " ", s).strip()


def set_para(p, text):
    if not p.runs:
        r = p.add_run(); r.text = text; return
    p.runs[0].text = text
    for extra in list(p.runs[1:]):
        extra._r.getparent().remove(extra._r)


def replace(slide_idx, pairs):
    slide = prs.slides[slide_idx]
    targets = {norm(o): n for o, n in pairs}
    done = set()
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            key = norm(p.text)
            if key in targets and key not in done:
                set_para(p, targets[key])
                done.add(key)
    missing = [o for o, _ in pairs if norm(o) not in done]
    for m in missing:
        print(f"  !! slide {slide_idx+1}: NOT FOUND -> {m!r}")
    return len(pairs) - len(missing)


# S6 — The Premise
replace(5, [
    ("This course teaches you to build with AI for AI",
     "This course teaches you to evaluate AI — to judge how much any number is worth."),
    ("Focus on designing systems that use LLMs as components: agents, tools, memory, orchestration.",
     "It is about reading evidence: benchmarks, demos, dashboards, vendor pitches, a teammate’s “it improved.”"),
    ("Build first, theorize after. Every concept tied to working code you write and evaluate.",
     "No code and no stats prerequisite. Every concept tied to an interactive playground you can feel."),
])

# S10 — What You Will Learn
replace(9, [
    ("Build reliable AI systems at scale", "Read any AI number as an estimate"),
    ("Use AI to build AI", "Build a test set you can trust"),
    ("Iterate towards success", "Tell a real improvement from noise"),
    ("Identify and master the key abstractions", "Report uncertainty honestly"),
    # 'Separate signal from hype' kept as-is
    ("Think in systems", "Ask the question that exposes a weak eval"),
])

# S11 — Course Philosophy
replace(10, [
    ("Hands-on. Become builders. Put something on your CV that actually matters",
     "Hands-on. Feel it, don’t just hear it."),
    ("Every lesson leads to working code.",
     "Every lesson has an interactive playground — no install, no code."),
    ("Systems thinking over prompt engineering", "Become hard to fool"),
    ("Prompts matter, but architecture matters more.",
     "The aim isn’t to learn statistics — it’s calibrated judgment about AI evidence."),
    # 'Evaluation before optimization' / "You can't improve what you can't measure." kept
    ("Understand and Embrace uncertainty", "Understand, report, and reduce uncertainty"),
])

# S12 — Course Structure
replace(11, [
    ("12 weeks · 24 lessons · 4 phases",
     "8 lessons · one source of uncertainty each · understand · report · reduce"),
    ("Phase 1", "Part 1"), ("Building AI Agents", "Foundations"),
    ("Weeks 1–2", "Lessons 1–3"),
    ("LLM APIs, context & memory,", "What “eval” is; assessing one run;"),
    ("agentic loops, tool integration", "proportions, probabilities, beliefs"),
    ("Phase 2", "Part 2"), ("Iterating Towards Success", "Overfitting & Judges"),
    ("Weeks 2–3", "Lessons 4–6"),
    ("Monitoring, quality estimation,", "Why re-running biases results;"),
    ("uncertainty, systematic iteration", "subjective judges, noisy truth"),
    ("Phase 3", "Part 3"), ("Complex Systems at Scale", "Compounding"),
    ("Weeks 4–6", "Lesson 7"),
    ("Multi-agent, design patterns,", "Drift and extra uncertainties —"),
    ("cost controls, frameworks", "and why they stack up"),
    ("Phase 4", "Part 4"), ("Capstone Project", "Capstone"),
    ("Weeks 6–12", "Lesson 8"),
    ("Spec, reviews, implementation,", "Design a reliable experiment,"),
    ("evaluation, final presentations", "end to end: question → decision"),
])

# S13 — Part 1
replace(12, [
    ("PHASE 1", "PART 1"), ("Building AI Agents", "Foundations"),
    ("Lessons 1–3  ·  Weeks 1–2", "Lessons 1–3"),
    ("L1: Hello World", "L1: “Eval” is now an Experiment"),
    ("Connecting to AI via API — chat, streaming, voice, image. Your first LLM call with latency measurement.",
     "From a check to an experiment. Today’s uncertainty: sampling noise — a score is one draw."),
    ("L2: Stateful Clients", "L2: Assessing one run of an agent"),
    ("Managing context and session memory. From stateless to stateful to memory-optimized agents. Cost and latency tradeoffs.",
     "Programmatic checks, AI judges, judges of judges. Uncertainty: variance around a score."),
    ("L3: Tools & The Agentic Loop", "L3: Proportions, Probabilities, Beliefs"),
    ("Tool calling, orchestration, the agent loop. Key principle: agents decide, tools execute.",
     "A gentle intro to probability — what are we actually measuring? Uncertainty: variability across customers and domains."),
])

# S14 — Part 2
replace(13, [
    ("PHASE 2", "PART 2"), ("Iterating Towards Success", "Overfitting & Judges"),
    ("Lessons 4–6  ·  Weeks 2–3", "Lessons 4–6"),
    ("L4: Monitoring & Business Assertions", "L4: Closed-Book Overfitting"),
    ("Observability, runtime checks, CI-integrated evaluation harnesses.",
     "Why “regression testing” gives biased results. Uncertainty: bias from multiple hypothesis testing."),
    ("L5: Quality & Performance Estimation", "L5: Open-Book Overfitting"),
    ("Metrics, golden datasets, LLM-as-judge. 'Optimizing in the Dark' — uncertainty quantification.",
     "Overfitting to the prompt, the benchmark, the eval set — and to the CTO. Spot it before it ships."),
    ("L6: Iterating on Clients, Agents & Prompts", "L6: Subjectivity & Noise in Judges"),
    ("Systematic improvement cycles, experiment design, error analysis.",
     "When the “right answer” is contested and the judge has its own opinions. Uncertainty: judge errors."),
])

# S15 — Part 3 (+ tools + capstone)
replace(14, [
    ("PHASE 3", "PART 3"), ("Complex Systems at Scale", "Compounding, Capstone & Tools"),
    ("Lessons 7–11  ·  Weeks 4–6", "Lessons 7–8 + course-wide tools"),
    ("Councils of Agents", "L7: Compounding Errors"),
    ("Diversity of opinions, ensemble approaches, consensus",
     "Drift and extra sources of error — and why they compound."),
    ("Building Complex Systems", "L8: Designing Reliable Experiments"),
    ("Divide et impera, testability, when NOT to use agents",
     "From a product question to a ship / don’t-ship decision."),
    ("Key Abstractions & Interfaces", "Ask the Playground"),
    ("Cost controls, guardrails, MCP protocol, observability",
     "A conversational explain-and-interrogate layer on every page."),
    ("Structuring Complex Projects", "Experiment-Report Evaluator"),
    ("Shipping into existing products, maintaining AI systems",
     "Submit an experiment write-up, get structured feedback."),
    ("AI Frameworks", "Capstone"),
    ("When and why to use them — LangChain case study",
     "Put the sources back together in one end-to-end eval."),
])

# S16 — Logistics & Policies
replace(15, [
    ("Basic software engineering experience preferred — but not required. Main prerequisite: passion to build great systems.",
     "No coding required. No probability or statistics background required — we demystify them through interactive experiments."),
    ("Expected for all lessons. Phase 4 sessions are critical — project reviews, working sessions, presentations.",
     "Expected for all lessons — each one builds a new source of uncertainty onto the last."),
    ("You must use AI tools to assist with coding and learning. But you must understand and explain all concepts in your submissions.",
     "Use AI tools freely to assist. But you must understand and be able to explain every number you report."),
])

# S17 — Closing
replace(16, [
    ("Let's Build.", "Let’s Measure."),
    ("Build first. Measure always. Ship with confidence.",
     "Don’t take the number at face value. Understand it, report it, reduce it."),
])

prs.save(SRC)
print("saved", SRC)
