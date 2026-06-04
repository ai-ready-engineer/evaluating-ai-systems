"""Append the Lesson-1 teaching slides to L1_outline.pptx.

Append-only: it adds the L1 section to whatever opening deck is already in L1_outline.pptx.
Run it through build_deck.py, which starts from the pristine opening every time so there are
no orphaned slide parts.

Slides follow the project slide conventions (see CLAUDE.md): every text run is >= 22pt,
no decorative panels/cards/colour-bars, and each slide is just a title, a short body, and
at most one chart — so they are easy to edit by hand.

    python build_deck.py        # full clean rebuild (recommended)
    python build_l1_slides.py   # append onto the current L1_outline.pptx
"""
from pathlib import Path
from PIL import Image

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = Path(__file__).resolve().parent
LESSON = HERE.parent
SRC = LESSON / "L1_generated.pptx"
CHARTS = LESSON / "assets" / "charts"

# Light mode: white background, dark text. (Charts are still dark PNGs for now.)
BG    = RGBColor(0xFF, 0xFF, 0xFF)
TEAL  = RGBColor(0x0F, 0x76, 0x6E)
TXT   = RGBColor(0x1A, 0x1F, 0x2E)
ICE   = RGBColor(0x47, 0x55, 0x69)
HEAD  = "Aptos Display"
BODY  = "Aptos"

MIN_PT = 22

prs = Presentation(str(SRC))
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[0]


def IN(v):
    return Emu(int(v * 914400))


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background(); r.shadow.inherit = False
    return s


def textbox(s, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(IN(l), IN(t), IN(w), IN(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, (text, size, color, bold) in enumerate(lines):
        assert size >= MIN_PT, f"font {size} below {MIN_PT}pt: {text[:40]}"
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(12); p.space_before = Pt(0)
        font = HEAD if bold and size >= 28 else BODY
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = font; r.font.color.rgb = color
    return tb


def title(s, text, size=32):
    textbox(s, 0.6, 0.45, 8.8, 1.2, [(text, size, TXT, True)])


def chart(s, name, top=1.55, bh=2.5):
    path = str(CHARTS / name)
    iw, ih = Image.open(path).size
    bl, bw = 0.6, 8.8
    ar = iw / ih
    w = bw; h = bw / ar
    if h > bh:
        h = bh; w = bh * ar
    l = bl + (bw - w) / 2
    s.shapes.add_picture(path, IN(l), IN(top), IN(w), IN(h))
    return top + h


def caption(s, text, top):
    textbox(s, 0.6, top + 0.15, 8.8, 1.2, [(text, 22, ICE, False)])


def code_lines(s, lines, top=1.6, size=22):
    tb = s.shapes.add_textbox(IN(0.7), IN(top), IN(8.6), IN(3.0))
    tf = tb.text_frame; tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, line in enumerate(lines):
        assert size >= MIN_PT, f"code font {size} below {MIN_PT}pt"
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(2); p.space_before = Pt(0)
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.name = "Courier New"; r.font.color.rgb = TXT
    return top + len(lines) * 0.38 + 0.12


# ---- append L1 section ----
print(f"appending L1 section onto {len(prs.slides._sldIdLst)} opening slides")

# 1 — divider
s = slide()
title(s, "Lesson 1 — Eval is now an Experiment", size=36)
textbox(s, 0.6, 2.1, 8.8, 3.0, [
    ("From software to ML to AI agents: an eval estimates a quantity — it doesn't just check a box.", 24, ICE, False),
    ("Today's source of uncertainty: sampling noise — a score is one draw.", 24, TEAL, True),
])

# 2 — one word, three meanings
s = slide()
title(s, "What does “eval” even mean?")
textbox(s, 0.6, 1.7, 8.8, 3.6, [
    ("Software 1.0 — a check. The suite passes or fails; a red test is a real regression.", 23, TXT, False),
    ("Machine learning — an experiment. A test set gives a score, and that score can move next run.", 23, TXT, False),
    ("AI agents — an open question. Many steps, judges, drift; even “correct” is part of the work.", 23, TXT, False),
])

# 2b — Software 1.0 concrete example
s = slide()
title(s, "Software 1.0 — a check you can test")
b = code_lines(s, [
    "def deviation_score(cpu_dev, mem_dev):",
    "    # each input in [0, 1]; 0 = at baseline",
    "    if not 0 <= cpu_dev <= 1:",
    "        raise ValueError(\"cpu_dev\")",
    "    if not 0 <= mem_dev <= 1:",
    "        raise ValueError(\"mem_dev\")",
    "    return (cpu_dev**2 + mem_dev**2) ** 0.5",
], top=1.5)
caption(s, "deviation_score(0.6, 0.8) == 1.0 — exact and reproducible; a failing test is a real regression.", b)

# 3 — the claim
s = slide()
title(s, "An eval is an experiment that estimates a quantity")
textbox(s, 0.6, 1.9, 8.8, 1.1, [("Model M has accuracy A on dataset D.", 32, TEAL, True)])
textbox(s, 0.6, 3.1, 8.8, 2.2, [
    ("Everything else — that A is reliable, holds on tomorrow's data, beats last quarter — sits on top of that one sentence.", 23, TXT, False),
    ("Today we take the first crack: even the honest A is just one draw.", 23, ICE, True),
])

# 4 — metric + confusion
s = slide()
title(s, "Pick a metric — and read it per class")
b = chart(s, "05_confusion.png")
caption(s, "Five categories, two classifiers, one 5×5 confusion matrix. The diagonal is what's right; zoom into any class for its precision and recall.", b)

# 5 — per-class story
s = slide()
title(s, "One number hides a per-class story")
b = chart(s, "04_perclass_f1.png")
caption(s, "Overall accuracy looks fine — but per class it is uneven. Always look under the headline.", b)

# 6 — jaggedness
s = slide()
title(s, "The jaggedness is real")
b = chart(s, "01_jaggedness.png")
caption(s, "Each dot is a ticket, coloured right/wrong. Right and wrong interlock — no clean boundary — and that doesn't vanish with more data.", b)

# 7 — the dark space
s = slide()
title(s, "We only see the points we labelled")
b = chart(s, "09_surface_reveal.png", bh=2.7)
caption(s, "We imagine a smooth competence surface; really we light only our test points and infer the rest of the room.", b)

# 8 — one draw / wobble
s = slide()
title(s, "Re-draw the test set and the score moves")
b = chart(s, "02_wobble.png")
caption(s, "Same fixed model, 2,000 random test sets of 50 — only the draw changed. 78% is one draw: quote n, and a range.", b)

# 9 — beta / reporting
s = slide()
title(s, "From one result to a credible range")
b = chart(s, "03_beta.png")
caption(s, "Report n, the 95% credible interval, and P(accuracy > X). The interval narrows as n grows — that is the value of a bigger test set.", b)

# 10 — class imbalance
s = slide()
title(s, "Accuracy can flatter a do-nothing model")
b = chart(s, "06_imbalance.png")
caption(s, "Skew the test set so one class dominates. A majority-only baseline gains accuracy while its macro-recall stays near 1/5 — a preview of L4.", b)

# 11 — trap + map
s = slide()
title(s, "A trap, and the map ahead")
textbox(s, 0.6, 1.7, 8.8, 3.6, [
    ("Trap: re-run the same eval and a few checks “fail” — that can be noise, not a regression. Re-run before believing; keep a sealed test set.", 23, TXT, False),
    ("Ahead: sampling noise (today), sampling bias, multiple comparisons, variance across domains, choice of metric, drift, judge noise, overfitting.", 23, ICE, False),
])

# 12 — takeaways
s = slide()
title(s, "Understand it · Report it · Reduce it")
textbox(s, 0.6, 1.7, 8.8, 3.6, [
    ("Understand — a score is one draw; sample size drives how much it wobbles.", 24, TEAL, True),
    ("Report — quote n, the 95% credible interval, and P(accuracy > X).", 24, ICE, True),
    ("Reduce — larger, fairer test set; re-run before believing; keep a sealed set.", 24, TXT, True),
])

prs.save(str(SRC))
print(f"saved {SRC.name} — now {len(prs.slides._sldIdLst)} slides")
