"""Capture L1_main.pptx into the reproducible build inputs used by build_l1_full.py.

Run this AFTER you have hand-edited L1_main.pptx in PowerPoint, to refresh:

  * lesson_1/assets/l1_spec.json     — per-slide text (runs: text/size/bold) + shape geometry
  * lesson_1/assets/deck/            — images embedded in the deck (per shape)
  * lesson_1/assets/deck_slides/     — full-slide PNG renders (faithful fallback art)

    python extract_l1_spec.py

It only READS L1_main.pptx. Slide rendering needs LibreOffice (`soffice`) + `pdftoppm`.
"""
import json
import shutil
import subprocess
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

HERE = Path(__file__).resolve().parent
LESSON = HERE.parent
MAIN = LESSON / "L1_main.pptx"
ASSETS = LESSON / "assets"
DECK = ASSETS / "deck"
SLIDES = ASSETS / "deck_slides"


def emu_in(v):
    try:
        return round(Emu(v).inches, 3)
    except Exception:
        return None


def extract_spec():
    DECK.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(MAIN))
    spec = []
    imgn = 0
    for i, s in enumerate(prs.slides):
        sl = {"n": i + 1, "shapes": []}

        def walk(shapes):
            nonlocal imgn
            for sh in shapes:
                d = {"type": str(sh.shape_type), "l": emu_in(sh.left), "t": emu_in(sh.top),
                     "w": emu_in(sh.width), "h": emu_in(sh.height)}
                if sh.shape_type == 13:
                    try:
                        imgn += 1
                        fn = f"s{i+1:02d}_{imgn}.{sh.image.ext}"
                        (DECK / fn).write_bytes(sh.image.blob)
                        d["img"] = fn
                    except Exception as e:
                        d["img"] = f"ERR {e}"
                if sh.has_text_frame and sh.text_frame.text.strip():
                    paras = []
                    for p in sh.text_frame.paragraphs:
                        runs = [{"t": r.text, "sz": (r.font.size.pt if r.font.size else None),
                                 "b": r.font.bold} for r in p.runs if r.text]
                        if runs:
                            paras.append(runs)
                    if paras:
                        d["paras"] = paras
                if getattr(sh, "has_table", False) and sh.has_table:
                    d["table"] = [[c.text for c in row.cells] for row in sh.table.rows]
                if sh.shape_type == 6:
                    d["group"] = True
                sl["shapes"].append(d)
                if sh.shape_type == 6:
                    walk(sh.shapes)

        walk(s.shapes)
        spec.append(sl)
    (ASSETS / "l1_spec.json").write_text(json.dumps(spec, indent=1, ensure_ascii=False))
    print(f"spec: {len(spec)} slides, {imgn} images -> {ASSETS/'l1_spec.json'}")


def render_slides():
    if not shutil.which("soffice"):
        print("soffice not found — skipping slide renders (keeping existing deck_slides/)")
        return
    SLIDES.mkdir(parents=True, exist_ok=True)
    subprocess.run(["soffice", "--headless", "--convert-to", "pdf", "--outdir", str(SLIDES), str(MAIN)],
                   check=True)
    pdf = SLIDES / "L1_main.pdf"
    subprocess.run(["pdftoppm", "-png", "-r", "150", str(pdf), str(SLIDES / "slide")], check=True)
    pdf.unlink(missing_ok=True)
    print(f"rendered slides -> {SLIDES}/slide-NN.png")


if __name__ == "__main__":
    extract_spec()
    render_slides()
