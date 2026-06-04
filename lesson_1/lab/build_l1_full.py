"""Rebuild the FULL Lesson-1 deck (L1_main's narrative) from scratch, in code.

Why this exists
---------------
L1_main.pptx was restructured by hand into a 49-slide narrative. This script makes that
deck reproducible:

  * Text-based slides are reconstructed as real, editable text + image shapes, driven by
    ``lesson_1/assets/l1_spec.json`` (extracted from the deck). Edit the copy there — or the
    per-slide overrides below — and re-run to regenerate.
  * The genuinely vector/diagram/table/card slides (role cards, function diagrams, scorecard
    tables, the numbered "What you'll learn / Philosophy / Structure" cards) are hand-designed
    in PowerPoint and don't survive clean code reconstruction, so they are placed as faithful
    rendered images from ``lesson_1/assets/deck_slides/``.

    python build_l1_full.py        # writes lesson_1/L1_generated.pptx

It never touches L1_main.pptx.

To re-capture spec + images after editing the deck by hand, see ``extract_l1_spec.py``.
"""
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).resolve().parent
LESSON = HERE.parent
ASSETS = LESSON / "assets"
SPEC = json.loads((ASSETS / "l1_spec.json").read_text())
DECK_IMG = ASSETS / "deck"            # images extracted from the deck (per shape)
SLIDE_IMG = ASSETS / "deck_slides"    # full-slide renders (faithful fallback)
OUT = LESSON / "L1_generated.pptx"

BG = RGBColor(0xFF, 0xFF, 0xFF)
TXT = RGBColor(0x1A, 0x1F, 0x2E)
BODY = "Aptos"
CHARTS = ASSETS / "charts"

# Hand-designed slides (title/cards/diagrams/tables/connectors, coloured badges, centred
# dividers, numbered lists) — these don't survive clean code reconstruction, so they are
# placed as faithful full-slide renders from assets/deck_slides/.
GRAPHIC = {1, 2, 3, 4, 5, 7, 9, 10, 11, 12, 14, 15, 21, 22, 23, 25, 26, 27, 32, 33,
           37, 38, 39, 40, 41, 42}

# The 7 chart slides are the original build_l1_slides.py charts (their embedded copies in
# L1_main have broken relationships) — regenerate them from the chart assets at slide geometry.
CHART_FALLBACK = {
    43: "05_confusion.png", 44: "04_perclass_f1.png", 45: "01_jaggedness.png",
    46: "09_surface_reveal.png", 47: "02_wobble.png", 48: "03_beta.png", 49: "06_imbalance.png",
}

EMU = 914400
def IN(v):
    return Emu(int(round(v * EMU)))

prs = Presentation()
prs.slide_width = IN(10.0)
prs.slide_height = IN(5.625)
BLANK = prs.slide_layouts[6]   # truly blank
W, H = prs.slide_width, prs.slide_height


def full_image(slide, path):
    slide.shapes.add_picture(str(path), 0, 0, W, H)


def add_text(slide, sh):
    tb = slide.shapes.add_textbox(IN(sh["l"]), IN(sh["t"]), IN(sh["w"]), IN(sh["h"]))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, runs in enumerate(sh["paras"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        for r in runs:
            run = p.add_run(); run.text = r["t"]
            run.font.name = BODY
            run.font.size = Pt(r["sz"] if r.get("sz") else 22)
            run.font.bold = bool(r.get("b"))
            run.font.color.rgb = TXT
    return tb


def build():
    for sl in SPEC:
        n = sl["n"]
        s = prs.slides.add_slide(BLANK)
        # white background
        from pptx.enum.shapes import MSO_SHAPE
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
        bg.fill.solid(); bg.fill.fore_color.rgb = BG
        bg.line.fill.background(); bg.shadow.inherit = False

        if n in GRAPHIC:
            full_image(s, SLIDE_IMG / f"slide-{n:02d}.png")
            continue

        for sh in sl["shapes"]:
            if "img" in sh:
                src = None
                if not str(sh["img"]).startswith("ERR"):
                    cand = DECK_IMG / sh["img"]
                    if cand.exists():
                        src = cand
                elif n in CHART_FALLBACK:
                    src = CHARTS / CHART_FALLBACK[n]
                if src:
                    s.shapes.add_picture(str(src), IN(sh["l"]), IN(sh["t"]), IN(sh["w"]), IN(sh["h"]))
            elif sh.get("paras"):
                add_text(s, sh)

    prs.save(str(OUT))
    print(f"built {OUT.name}: {len(prs.slides._sldIdLst)} slides "
          f"({len(GRAPHIC)} as images, {len(SPEC) - len(GRAPHIC)} reconstructed)")


if __name__ == "__main__":
    build()
